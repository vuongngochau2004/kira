"""Document extraction using PyMuPDF with PaddleOCR fallback."""

import sys
import logging
from pathlib import Path
from dataclasses import dataclass, field

# Add parent directory to path for imports
sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

logger = logging.getLogger(__name__)


@dataclass
class ExtractionResult:
    """Result from text extraction."""
    text: str = ""
    pages: int = 0
    metadata: dict = field(default_factory=dict)
    success: bool = True
    error: str | None = None


SUPPORTED_TYPES = {
    "pdf", "docx", "pptx", "xlsx", "html", "md", "txt",
}

OCR_TYPES = {"png", "jpg", "jpeg", "tiff", "bmp", "gif", "webp"}


async def extract_pdf(file_path: str, use_ocr_fallback: bool = True) -> ExtractionResult:
    """Extract text from PDF using PyMuPDF with OCR fallback for scanned pages.

    Args:
        file_path: Path to PDF file
        use_ocr_fallback: Enable OCR for pages with no extractable text

    Returns:
        ExtractionResult with text and metadata
    """
    logger.info("Extracting PDF using PyMuPDF (native text): %s", file_path)
    try:
        import fitz
        from src.ingestion.paddleocr_client import get_ocr_client

        doc = fitz.open(file_path)
        num_pages = len(doc)
        text_parts = []
        pages_needing_ocr = []

        # First pass: extract text with PyMuPDF
        for page_num, page in enumerate(doc):
            page_text = page.get_text()
            if page_text.strip():
                text_parts.append(page_text)
            else:
                # Mark page for OCR - use placeholder
                text_parts.append("")  # Placeholder for OCR result
                pages_needing_ocr.append(page_num)

        # If all pages have text or OCR disabled, return early
        if not pages_needing_ocr or not use_ocr_fallback:
            doc.close()
            full_text = "\n\n".join([t for t in text_parts if t])
            return ExtractionResult(
                text=full_text,
                pages=num_pages,
                metadata={"extractor": "pymupdf", "engine": "fitz", "ocr_used": False},
                success=bool(full_text.strip()),
            )

        # Second pass: OCR for pages without text
        ocr_client = get_ocr_client()
        async with ocr_client:
            for page_num in pages_needing_ocr:
                page = doc[page_num]

                # Convert page to image
                mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better OCR
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")

                # Perform OCR
                ocr_result = await ocr_client.ocr_image_bytes(img_bytes)

                if ocr_result.success and ocr_result.text.strip():
                    # Update placeholder with OCR text
                    text_parts[page_num] = ocr_result.text

        doc.close()

        full_text = "\n\n".join([t for t in text_parts if t])

        return ExtractionResult(
            text=full_text,
            pages=num_pages,
            metadata={
                "extractor": "pymupdf-ocr-hybrid",
                "engine": "fitz",
                "ocr_used": True,
                "ocr_pages": len(pages_needing_ocr),
            },
            success=bool(full_text.strip()),
        )

    except Exception as e:
        logger.error("PyMuPDF extraction failed for %s: %s", file_path, e)
        return ExtractionResult(
            success=False,
            error=f"PDF extraction failed: {e}",
        )


def extract_docx(file_path: str) -> ExtractionResult:
    """Extract text from DOCX using python-docx.

    Args:
        file_path: Path to DOCX file

    Returns:
        ExtractionResult with text and metadata
    """
    logger.info("Extracting DOCX using python-docx: %s", file_path)
    try:
        from docx import Document

        doc = Document(file_path)
        text_parts = []

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        full_text = "\n".join(text_parts)

        return ExtractionResult(
            text=full_text,
            pages=1,
            metadata={"extractor": "python-docx", "extraction_method": "normal", "ocr_used": False},
            success=bool(full_text.strip()),
        )
    except Exception as e:
        logger.error("DOCX extraction failed for %s: %s", file_path, e)
        return ExtractionResult(
            success=False,
            error=f"DOCX extraction failed: {e}",
        )


def extract_pptx(file_path: str) -> ExtractionResult:
    """Extract text from PPTX using python-pptx.

    Args:
        file_path: Path to PPTX file

    Returns:
        ExtractionResult with text and metadata
    """
    logger.info("Extracting PPTX using python-pptx: %s", file_path)
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts = []

        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append("\n".join(slide_text))

        full_text = "\n\n".join(text_parts)

        return ExtractionResult(
            text=full_text,
            pages=len(prs.slides),
            metadata={"extractor": "python-pptx", "extraction_method": "normal", "ocr_used": False},
            success=bool(full_text.strip()),
        )
    except Exception as e:
        logger.error("PPTX extraction failed for %s: %s", file_path, e)
        return ExtractionResult(
            success=False,
            error=f"PPTX extraction failed: {e}",
        )


def extract_text_file(file_path: str) -> ExtractionResult:
    """Extract text from plain text or markdown file.

    Args:
        file_path: Path to text file

    Returns:
        ExtractionResult with text and metadata
    """
    logger.info("Extracting text file: %s", file_path)
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()

        return ExtractionResult(
            text=text,
            pages=1,
            metadata={"extractor": "native", "extraction_method": "normal", "ocr_used": False},
            success=bool(text.strip()),
        )
    except UnicodeDecodeError:
        # Try with different encoding
        logger.warning("UTF-8 decoding failed for %s, falling back to latin-1", file_path)
        try:
            with open(file_path, "r", encoding="latin-1") as f:
                text = f.read()
            return ExtractionResult(
                text=text,
                pages=1,
                metadata={"extractor": "native", "encoding": "latin-1", "extraction_method": "normal", "ocr_used": False},
                success=True,
            )
        except Exception as e:
            logger.error("Text file extraction failed for %s with latin-1: %s", file_path, e)
            return ExtractionResult(
                success=False,
                error=f"Text extraction failed: {e}",
            )
    except Exception as e:
        logger.error("Text file extraction failed for %s: %s", file_path, e)
        return ExtractionResult(
            success=False,
            error=f"Text extraction failed: {e}",
        )


async def extract_image_ocr(file_path: str) -> ExtractionResult:
    """Extract text from image using PaddleOCR service.

    Args:
        file_path: Path to image file

    Returns:
        ExtractionResult with text and metadata
    """
    try:
        from src.ingestion.paddleocr_client import get_ocr_client

        ocr_client = get_ocr_client()
        async with ocr_client:
            result = await ocr_client.ocr_file(file_path)

        if result.success:
            return ExtractionResult(
                text=result.text,
                pages=1,
                metadata={
                    "extractor": "paddleocr",
                    "confidence": result.confidence,
                },
                success=True,
            )
        else:
            # Fallback to Docling if PaddleOCR fails
            return await _extract_image_docling(file_path)

    except Exception as e:
        # Fallback to Docling on any error
        return await _extract_image_docling(file_path, error=str(e))





async def _extract_image_docling(file_path: str, error: str | None = None) -> ExtractionResult:
    """Fallback image extraction using Docling.

    Args:
        file_path: Path to image file
        error: Original error that triggered fallback

    Returns:
        ExtractionResult with text and metadata
    """
    logger.info("Extracting image using Docling OCR: %s", file_path)
    try:
        from docling.document_converter import DocumentConverter

        converter = DocumentConverter()
        result = converter.convert(file_path)

        text = result.document.export_to_markdown()

        return ExtractionResult(
            text=text,
            pages=1,
            metadata={"extractor": "docling-ocr", "fallback_from": "paddleocr"},
            success=bool(text.strip()),
        )
    except Exception as e:
        fallback_err = str(e)
        return ExtractionResult(
            success=False,
            error=f"Image OCR failed (PaddleOCR: {error}, Docling: {fallback_err})",
        )


async def extract_content(file_path: str, file_type: str) -> ExtractionResult:
    """Extract text from file using appropriate extractor.

    Args:
        file_path: Path to the file
        file_type: File extension (e.g., "pdf", "docx")

    Returns:
        ExtractionResult with text, pages, metadata
    """
    from pathlib import Path

    if not Path(file_path).exists():
        return ExtractionResult(
            success=False,
            error=f"File not found: {file_path}",
        )

    file_type = file_type.lower().strip().lstrip(".")

    if file_type not in SUPPORTED_TYPES and file_type not in OCR_TYPES:
        return ExtractionResult(
            success=False,
            error=f"Unsupported file type: {file_type}",
        )

    is_ocr = file_type in OCR_TYPES
    method_name = "OCR" if is_ocr else "normal (text-based)"
    print(f"[Extractor] Starting {method_name} extraction for file: {file_path} (format: {file_type})")
    logger.info("Starting %s extraction for file: %s (format: %s)", method_name, file_path, file_type)

    try:
        if file_type == "pdf":
            return await extract_pdf(file_path)
        elif file_type == "docx":
            result = extract_docx(file_path)
        elif file_type == "pptx":
            result = extract_pptx(file_path)
        elif file_type in {"txt", "md"}:
            result = extract_text_file(file_path)
        elif file_type in OCR_TYPES:
            return await extract_image_ocr(file_path)
        else:
            result = ExtractionResult(
                success=False,
                error=f"Format {file_type} not yet supported",
            )

        if result.success:
            if "extraction_method" not in result.metadata:
                result.metadata["extraction_method"] = "ocr" if is_ocr else "normal"
            if "ocr_used" not in result.metadata:
                result.metadata["ocr_used"] = is_ocr
            print(f"[Extractor] Extraction successful. Method: {result.metadata['extraction_method'].upper()}, Pages: {result.pages}")
            logger.info("Extraction successful. Method: %s, Pages: %d", result.metadata['extraction_method'], result.pages)
        else:
            print(f"[Extractor ERROR] Extraction failed: {result.error}")
            logger.error("Extraction failed: %s", result.error)

        return result

    except Exception as e:
        logger.error("Extraction failed for %s: %s", file_path, e)
        return ExtractionResult(
            success=False,
            error=f"Extraction failed: {e}",
        )


def extract_content_sync(file_path: str, file_type: str) -> ExtractionResult:
    """Synchronous wrapper for extract_content.

    Args:
        file_path: Path to the file
        file_type: File extension (e.g., "pdf", "docx")

    Returns:
        ExtractionResult with text, pages, metadata
    """
    import asyncio

    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            # We're in an async context, use create_task
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor() as pool:
                future = pool.submit(
                    asyncio.run,
                    extract_content(file_path, file_type)
                )
                return future.result()
        else:
            # No event loop, run directly
            return asyncio.run(extract_content(file_path, file_type))
    except RuntimeError:
        # No event loop, create new one
        return asyncio.run(extract_content(file_path, file_type))


__all__ = [
    "ExtractionResult",
    "extract_content",
    "extract_content_sync",
    "extract_pdf",
    "extract_docx",
    "extract_pptx",
    "extract_text_file",
    "extract_image_ocr",
    "_extract_image_docling",
]
